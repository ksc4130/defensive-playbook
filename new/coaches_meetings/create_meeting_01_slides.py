#!/usr/bin/env python3
"""
Generate Meeting 1 slide deck: Introduction to the Defense
River Valley Vikings | 2026 Season

Output: .pptx that can be uploaded directly to Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────────
COLUMBIA_BLUE = RGBColor(0x6C, 0xAC, 0xE4)
GOLD = RGBColor(0xCF, 0xA7, 0x00)
NAVY = RGBColor(0x00, 0x2D, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
FAINT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, bold_prefix=True, line_spacing=1.3,
                    font_name="Calibri"):
    """Add a bulleted list. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)

        if isinstance(item, tuple):
            # Bold prefix + normal rest
            run_b = p.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = color
            run_b.font.bold = True
            run_b.font.name = font_name

            run_n = p.add_run()
            run_n.text = item[1]
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
            run_n.font.bold = False
            run_n.font.name = font_name
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = font_name

    return txBox


def add_gold_bar(slide, top, width=SLIDE_W):
    """Thin gold accent bar."""
    add_shape(slide, Inches(0), top, width, Pt(4), GOLD)


def section_header_slide(prs, title, subtitle=None):
    """Navy background section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(2.8))
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
                 title, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1),
                     subtitle, font_size=22, color=COLUMBIA_BLUE, bold=False,
                     alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title):
    """White background content slide with navy top bar and gold accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_gold_bar(slide, Inches(1.1))
    # Title text
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.85),
                 title, font_size=32, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    # Footer
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(5), Inches(0.4),
                 "River Valley Vikings | 2026 Defense", font_size=10,
                 color=MED_GRAY, bold=False)
    return slide


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, header_color=NAVY, header_text_color=WHITE,
              font_size=14, row_alt_color=FAINT_BLUE):
    """Add a formatted table. data is list of lists (rows x cols)."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            # Cell fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ── Build Presentation ──────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Gold accent bars
    add_gold_bar(slide, Inches(1.8))
    add_gold_bar(slide, Inches(5.3))

    add_text_box(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(1.2),
                 "RIVER VALLEY VIKINGS", font_size=52, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.8),
                 "2026 DEFENSIVE SYSTEM", font_size=36, color=COLUMBIA_BLUE,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.8),
                 "Coaches Introduction  |  Meeting 1", font_size=22,
                 color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.7),
                 "Confidential \u2014 Coaching Staff Only", font_size=16,
                 color=MED_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Agenda
    # ════════════════════════════════════════════════════════════════════════
    slide = content_slide(prs, "TODAY\u2019S AGENDA")

    items = [
        ("1.  ", "Defensive Identity & Philosophy"),
        ("2.  ", "Position Designations \u2014 Our Language"),
        ("3.  ", "Call Grammar \u2014 How We Communicate"),
        ("4.  ", "The System Overview \u2014 What\u2019s in the Toolbox"),
        ("5.  ", "How It Fits Together"),
        ("6.  ", "What Comes Next"),
    ]
    add_bullet_list(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5),
                    items, font_size=26, line_spacing=1.6)

    add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6),
                 "Goal: Understand WHO we are, HOW we talk, and WHAT the system looks like.",
                 font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: DEFENSIVE IDENTITY
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "DEFENSIVE IDENTITY", "Who We Are")

    # ── Slide 4 — Philosophy Quote ────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(1.5))
    add_gold_bar(slide, Inches(5.8))

    quote = (
        "\u201CWe play fast, physical, and disciplined football.\n"
        "Our job is to remove explosive plays, control the box,\n"
        "and force the offense to execute perfectly\n"
        "for an entire drive.\u201D"
    )
    txBox = add_text_box(slide, Inches(1.2), Inches(2.0), Inches(10.9),
                         Inches(3.5), "", font_size=30, color=WHITE)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(44)
    run = p.add_run()
    run.text = quote
    run.font.size = Pt(28)
    run.font.color.rgb = WHITE
    run.font.italic = True
    run.font.name = "Calibri"

    # ── Slide 5 — Base Structure ──────────────────────────────────────────
    slide = content_slide(prs, "BASE STRUCTURE")

    data = [
        ["Element", "Detail"],
        ["Personnel", "4-2-5 base; 3-4 packages (Mint / Ace / Jet / Slip)"],
        ["Strength", 'Field ("protect space")'],
        ["Shell", "Two-high (pre-snap disguise)"],
        ["Box", "6-man"],
        ["Call Format", "Front + Stunt + Blitz + Coverage/Tags"],
        ["Communication", "Signaled (front/stunt/blitz); coverage is verbal"],
    ]
    add_table(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5),
              rows=7, cols=2, data=data,
              col_widths=[Inches(3), Inches(7.9)], font_size=18)

    add_text_box(slide, Inches(1.2), Inches(6.3), Inches(10.9), Inches(0.5),
                 "Multiple looks, simple rules. The picture changes \u2014 the principles don\u2019t.",
                 font_size=16, color=NAVY, bold=True)

    # ── Slide 6 — Non-Negotiables ─────────────────────────────────────────
    slide = content_slide(prs, "THE NON-NEGOTIABLES")

    items = [
        ("1. No explosives. ", "Make them snap it again."),
        ("2. Aggressive with discipline. ", "Attack without losing leverage or committing penalties."),
        ("3. Box the run. ", "Set hard edges, compress gaps, keep the ball in help."),
        ("4. Relentless pursuit. ", "11 to the ball, correct angles, swarming tackles."),
        ("5. Win leverage. ", "Outside vs perimeter; inside vs quick game."),
        ("6. Win situations. ", "1st down sets the series; 3rd down = get off the field."),
        ("7. Smart aggression. ", "Plus-one rushers to blockers with zero coverage behind it. No-win situation for the offense."),
        ("8. Communication standard. ", "Simple, loud, early. Everyone echoes."),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5),
                    items, font_size=18, line_spacing=1.35)

    # ── Slide 7 — Grading Standard ────────────────────────────────────────
    slide = content_slide(prs, "THE STANDARD")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
                 "Six standards. No exceptions.",
                 font_size=18, color=MED_GRAY)

    items = [
        ("No lost edges ", "\u2014 contain player keeps QB and ball inside"),
        ("No uncovered gaps ", "\u2014 adjust instantly if gap changes with stunt/pressure"),
        ("No free releases ", "\u2014 collision every route; reroute, disrupt timing, deny easy access"),
        ("No missed tackles ", "\u2014 leverage + near foot + wrap; eliminate YAC"),
        ("No penalties ", "that extend drives"),
        ("No loafs ", "\u2014 pursuit is mandatory"),
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.5),
                    items, font_size=22, line_spacing=1.5)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: POSITION DESIGNATIONS
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "POSITION DESIGNATIONS", "Our Language")

    # ── Slide 9 — Position Table ──────────────────────────────────────────
    slide = content_slide(prs, "THE 11 POSITIONS")

    data = [
        ["Letter", "Name", "Identity"],
        ["FC", "Field Corner", "Lockdown corner. Speed, cover skills, competitive."],
        ["FS", "Field Safety", "Quarterback of the secondary. Smart, vocal, range."],
        ["B", "Bandit (Nickel)", "Swiss army knife. Best athlete on defense. DB/OLB hybrid."],
        ["A", "Anchor (Field DE)", "Edge setter. Anchor strength, contain discipline."],
        ["T", "Tackle (Field DT)", "Interior disruptor. Size, anchor, pass-rush ability."],
        ["N", "Nose (Boundary DT)", "Interior disruptor. Size, anchor, pass-rush ability."],
        ["E", "Edge (Boundary DE)", "Boundary edge. Versatile \u2014 DE in 4-down, OLB in 3-down."],
        ["M", "Mike (Field ILB)", "Downhill runner. Instincts, tackling, field-side playmaker."],
        ["W", "Will (Boundary ILB)", "Cutback hunter. Quick, disciplined, boundary-side playmaker."],
        ["D", "Dawg (Boundary S)", "Boundary fixer. Physical, smart, force player."],
        ["BC", "Boundary Corner", "Lockdown corner. Speed, cover skills, competitive."],
    ]
    add_table(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.8),
              rows=12, cols=3, data=data,
              col_widths=[Inches(1.5), Inches(3.5), Inches(7.1)], font_size=16)

    # ── Slide 10 — Defaults & Key Roles ───────────────────────────────────
    slide = content_slide(prs, "DEFAULT SIDES & KEY ROLES")

    # Field side box
    box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2),
                    FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.5),
                 "FIELD SIDE (wide side)", font_size=22, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(1.4),
                    ["A (Anchor) \u2014 Field DE",
                     "T (Tackle) \u2014 Field DT",
                     "B (Bandit) \u2014 Nickel / Field force",
                     "FC / FS \u2014 Field corner & safety"],
                    font_size=16, color=DARK_GRAY)

    # Boundary side box
    box2 = add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.2),
                     FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.1), Inches(0.5),
                 "BOUNDARY SIDE (short side)", font_size=22, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(1.4),
                    ["E (Edge) \u2014 Boundary DE",
                     "N (Nose) \u2014 Boundary DT",
                     "D (Dawg) \u2014 Boundary safety / force",
                     "BC \u2014 Boundary corner"],
                    font_size=16, color=DARK_GRAY)

    # Key distinctions
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
                 "Key Distinctions", font_size=22, color=NAVY, bold=True)
    items = [
        ("Contain: ", "Maintain outside leverage. Squeeze run inside. Set the edge."),
        ("Force: ", "Collision at or behind LOS. Turn the ball inside. More aggressive than contain."),
        ("Bandit (B): ", "Most versatile player on the field. DB in coverage, OLB in 3-down packages."),
        ("Safeties are fixers. ", "FS owns field-side calls. D owns boundary-side calls."),
    ]
    add_bullet_list(slide, Inches(1.0), Inches(4.7), Inches(11.5), Inches(2.2),
                    items, font_size=16, line_spacing=1.4)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: TECHNIQUE & GAP TERMINOLOGY
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "TECHNIQUE & GAP NAMES", "Alignment Language")

    # ── DL Technique Chart ───────────────────────────────────────────────
    slide = content_slide(prs, "DL TECHNIQUE CHART")

    data = [
        ["Tech", "Alignment"],
        ["0", "Head up on Center"],
        ["1", "Shade on Center"],
        ["2i", "Inside shade on Guard"],
        ["2", "Head up on Guard"],
        ["3", "Outside shade on Guard"],
        ["4i", "Inside shade on Tackle"],
        ["4", "Head up on Tackle"],
        ["5", "Outside shade on Tackle"],
        ["6i", "Inside shade on TE"],
        ["6", "Head up on TE"],
        ["7", "Outside shade on TE"],
        ["8i", "Inside shade on Wing"],
        ["8", "Head up on Wing"],
        ["9", "Outside shade on Wing"],
    ]
    add_table(slide, Inches(1.5), Inches(1.35), Inches(10.3), Inches(5.8),
              rows=15, cols=2, data=data,
              col_widths=[Inches(2), Inches(8.3)], font_size=16)

    # ── Gap Names ────────────────────────────────────────────────────────
    slide = content_slide(prs, "GAP NAMES")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Every gap has a name. Every defender owns a gap.",
                 font_size=18, color=MED_GRAY)

    data = [
        ["Gap", "Location"],
        ["A Gap", "Between Center and Guard"],
        ["B Gap", "Between Guard and Tackle"],
        ["C Gap", "Between Tackle and TE (or outside Tackle if no TE)"],
        ["D Gap", "Outside the TE"],
        ["E Gap", "Outside the Wing"],
        ["Alley", "Outside contain lane \u2014 where the force / contain player operates"],
    ]
    add_table(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.8),
              rows=7, cols=2, data=data,
              col_widths=[Inches(2), Inches(8.3)], font_size=18)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
                 "Gaps exist on both sides of the center. A/B/C/D field + A/B/C/D boundary.",
                 font_size=16, color=NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: CALL GRAMMAR
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "CALL GRAMMAR", "How We Communicate")

    # ── Slide 12 — The Formula ────────────────────────────────────────────
    slide = content_slide(prs, "THE FORMULA")

    # Big formula display
    formula_box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7),
                            Inches(1.4), FAINT_BLUE, COLUMBIA_BLUE)
    # Build the formula with colored segments
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    parts = [
        ("FRONT", NAVY),
        ("  +  ", DARK_GRAY),
        ("STUNT", RGBColor(0x8B, 0x00, 0x00)),  # dark red
        ("  +  ", DARK_GRAY),
        ("BLITZ", RGBColor(0x00, 0x6B, 0x3C)),  # dark green
        ("  +  ", DARK_GRAY),
        ("COVERAGE / TAGS", GOLD),
    ]
    for text, color in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = "Calibri"

    # Examples
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
                 "Examples:", font_size=22, color=NAVY, bold=True)

    examples = [
        # Each example: list of (text, color) tuples
        [("SHADE", NAVY), (" + ", DARK_GRAY), ("NINJA", GOLD)],
        [("UNDER", NAVY), (" + ", DARK_GRAY), ("SLANT", RGBColor(0x8B, 0x00, 0x00)),
         (" + ", DARK_GRAY), ("ZEUS", GOLD)],
        [("GRIZZLY", NAVY), (" + ", DARK_GRAY), ("sWarM", RGBColor(0x00, 0x6B, 0x3C)),
         (" + ", DARK_GRAY), ("ZORRO", GOLD)],
        [("EYES", NAVY), (" + ", DARK_GRAY), ("JACKS", RGBColor(0x8B, 0x00, 0x00)),
         (" + ", DARK_GRAY), ("sWarM", RGBColor(0x00, 0x6B, 0x3C)),
         (" + ", DARK_GRAY), ("ZORRO", GOLD)],
    ]

    y_pos = Inches(3.9)
    for ex in examples:
        txBox = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(10), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        for text, color in ex:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(22)
            run.font.bold = True if text.strip() not in ["+"] else False
            run.font.color.rgb = color
            run.font.name = "Calibri"
        y_pos += Inches(0.55)

    # Bottom note
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
                 "Front/stunt/blitz calls signaled \u2014 coverage is verbal. Safeties echo and relay.",
                 font_size=16, color=MED_GRAY, bold=False)

    # ── Slide 13 — Tempo Call ─────────────────────────────────────────────
    slide = content_slide(prs, "TEMPO CALL")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8),
                 "When the offense\u2019s tempo is faster than we can signal our plays in,\n"
                 "the defense defaults to the tempo call.",
                 font_size=20, color=DARK_GRAY)

    # Big tempo call box
    box = add_shape(slide, Inches(3), Inches(2.8), Inches(7.3), Inches(1.5),
                    NAVY, GOLD)
    add_text_box(slide, Inches(3), Inches(2.9), Inches(7.3), Inches(1.3),
                 "DEFAULT:  SHADE + NINJA", font_size=40, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    items = [
        "Set before each possession. Can change at any point during the game.",
        "Safeties must know the current tempo call at all times.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(4.8), Inches(10.9), Inches(2),
                    items, font_size=18, line_spacing=1.5)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: SYSTEM OVERVIEW
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "THE SYSTEM", "What\u2019s in the Toolbox")

    # ── Slide 15 — Fronts Menu ────────────────────────────────────────────
    slide = content_slide(prs, "FRONT MENU  \u2014  4-DOWN")

    data = [
        ["Front", "Interior (A/T/N/E)", "When to Call"],
        ["SHADE", "5 / 3 / 2i / 5", "Base vs spread / RPO. Where we live."],
        ["UNDER", "5 / 2i / 3 / 5", "Tendency breaker. Boundary-run weeks."],
        ["EYES", "5 / 2i / 2i / 5", "Vs zone / duo. Square the interior."],
        ["WIDE", "5 / 3 / 3 / 5", "Force the bounce. Vs B-gap heavy / gap schemes."],
        ["DEUCES", "5 / 2 / 2 / 5", "Vs Wing-T / pullers."],
        ["GRIZZLY", "4i / 2i / 2i / 4i", "Goal line. Power / counter / zone. Edge pressure."],
        ["BOSS", "5 / 3 / 1 / 5", "Overload field / confuse OL. SPLIT \u2192 SHADE post-snap. ANGLE \u2192 EYES post-snap."],
        ["BOSS UNDER", "5 / 1 / 3 / 5", "Overload boundary / confuse OL. SPLIT or ANGLE combos."],
    ]
    add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5),
              rows=9, cols=3, data=data,
              col_widths=[Inches(2.2), Inches(2.5), Inches(7.8)], font_size=16)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
                 "TE SET rule overrides all fronts except GRIZZLY when TE is attached.\n"
                 "DISCO: Week-by-week SHADE adjustment based on the back and offensive tendencies.",
                 font_size=15, color=NAVY, bold=True)

    # ── Slide 16 — 3-Down Packages ────────────────────────────────────────
    slide = content_slide(prs, "3-DOWN PACKAGES")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Our odd front package. Keeps LBs clean, controls the center, adjusts to personnel.",
                 font_size=18, color=MED_GRAY)

    data = [
        ["Package", "Interior (A/T/N)", "Key Feature"],
        ["MINT", "4i / 0 / 4i", "Base 3-down. Contain adjusts by coverage."],
        ["ACE", "4 / 0 / 4", "2-gap interior. OLBs contain."],
        ["JET", "5 / 0 / 5", "Edges at 5. A and N are contain."],
        ["SLIP", "5 / 0 / 4i", "Penetrating. Takes B out of run/pass conflict."],
    ]
    add_table(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.2),
              rows=5, cols=3, data=data,
              col_widths=[Inches(2.5), Inches(2.5), Inches(6.7)], font_size=18)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
                 "B = Field OLB.  E = Boundary OLB.  Strength = Field.",
                 font_size=16, color=NAVY, bold=True)

    # ── Slide 17 — Stunts Overview ────────────────────────────────────────
    slide = content_slide(prs, "STUNT MENU")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Coordinated DL movement. Every stunt has rules for who moves and where.",
                 font_size=18, color=MED_GRAY)

    # Two columns — with 4-down / 3-down specifics
    col1 = [
        ("SLANT ", "\u2014 Slant to field. 4-down: T+N. 3-down: A+T+N."),
        ("ANGLE ", "\u2014 Slant to boundary. 4-down: T+N. 3-down: A+T+N."),
        ("PINCH ", "\u2014 Pinch inside. 4-down: T+N (A gaps). 3-down: A+N (B gaps)."),
        ("JACKS ", "\u2014 Shoot out. 4-down: T+N (B gaps). 3-down: A+N (C gaps)."),
        ("VEER ", "\u2014 Slant to motion / to TE. 4-down: T+N. 3-down: A+T+N."),
        ("SPLIT ", "\u2014 BOSS / BOSS UNDER / GRIZZLY only"),
    ]
    col2 = [
        ("CRASH ", "\u2014 Everyone shoots inside"),
        ("ANCHOR ATTACK ", "\u2014 A washes OL inside"),
        ("EDGE ATTACK ", "\u2014 E washes OL inside"),
        ("ANCHOR RAVEN ", "\u2014 A shoots B gap"),
        ("EDGE RAVEN ", "\u2014 E shoots B gap"),
    ]

    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4),
                    col1, font_size=16, line_spacing=1.5)
    add_bullet_list(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4),
                    col2, font_size=16, line_spacing=1.5)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
                 "Front-stunt legality matrix exists \u2014 not every stunt works with every front.",
                 font_size=15, color=NAVY, bold=True)

    # ── Slide 18 — Pressures Overview ─────────────────────────────────────
    slide = content_slide(prs, "PRESSURE MENU")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Pressure to take away something \u2014 not just to be aggressive.",
                 font_size=18, color=MED_GRAY)

    # Three column layout
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(0.5),
                 "Singles", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(3.5), Inches(2.5),
                    ["Mike", "Will", "Bandit", "Dawg"],
                    font_size=18, line_spacing=1.4)

    add_text_box(slide, Inches(4.8), Inches(2.2), Inches(3.5), Inches(0.5),
                 "Combos", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(4.8), Inches(2.7), Inches(3.5), Inches(2.5),
                    ["sWarM (M+W)", "BooM (B+M)", "BoW (B+W)",
                     "MaD (M+D)", "Eat (M+W+B)"],
                    font_size=18, line_spacing=1.4)

    add_text_box(slide, Inches(8.8), Inches(2.2), Inches(4), Inches(0.5),
                 "Packaged", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(8.8), Inches(2.7), Inches(4), Inches(3.5),
                    [("Hammer ", "\u2014 B edge + Anchor Attack"),
                     ("Shave ", "\u2014 W edge + Edge Attack"),
                     ("Bandit Raven ", "\u2014 B blitzes C gap, A shoots B gap"),
                     ("Will Raven ", "\u2014 W blitzes C gap, E shoots B gap")],
                    font_size=18, line_spacing=1.4)

    # MUG note
    box = add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2),
                    FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.0),
                 "MUG: M and/or W walk up to the LOS at their gap. Same job, different picture.\n"
                 "Weekly menu locked by Friday. Game-plan specials installed as needed.",
                 font_size=16, color=DARK_GRAY)

    # ── Slide 19 — Coverage Families ──────────────────────────────────────
    slide = content_slide(prs, "THE FOUR COVERAGE FAMILIES")

    # NINJA box
    y = Inches(1.5)
    box_h = Inches(1.15)
    gap = Inches(0.15)

    box = add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(10), Inches(0.5),
                 "NINJA  \u2014  Flip (SKY/MOD) / Clamp (Cloud/Read 2) / Poach", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Pattern-match quarters. Our base coverage. Most of our snaps.",
                 font_size=16, color=DARK_GRAY)

    y += box_h + gap
    box = add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 "COVER 1  (Man-Match)", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Oregon (FS = MOF) / Oklahoma (D = MOF) / Ohio (B = MOF). Man assignments + MOF safety help.",
                 font_size=16, color=DARK_GRAY)

    y += box_h + gap
    box = add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.5),
                 "COVER 0 / Z-FAMILY", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Zero coverage. Zeus / Zorro / Zunnel / Zill / Zike. Run-first, delayed pressure. No deep help.",
                 font_size=16, color=DARK_GRAY)

    y += box_h + gap
    box = add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 "VIKING  (Cover 3)", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Spot-drop Cover 3. Safe, anti-explosive.",
                 font_size=16, color=DARK_GRAY)

    # Bottom note
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
                 "Zeus is the pressure. Viking is the safety net. NINJA is the engine of the defense.",
                 font_size=16, color=NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: HOW IT FITS TOGETHER
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "HOW IT FITS TOGETHER",
                         "Same picture pre-snap. Different answers post-snap.")

    # ── Slide 21 — The Concept ────────────────────────────────────────────
    slide = content_slide(prs, "ONE LOOK, MULTIPLE ANSWERS")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8),
                 "We have a primary two-high shell and a three-high shell. Both have multiple\n"
                 "post-snap pictures. The offense can\u2019t tell what\u2019s coming.",
                 font_size=20, color=DARK_GRAY)

    # Three boxes showing the concept
    box_w = Inches(3.5)
    box_h = Inches(3.0)
    start_x = Inches(0.8)
    gap_x = Inches(0.55)

    # Box 1
    x = start_x
    add_shape(slide, x, Inches(2.8), box_w, box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, x + Inches(0.15), Inches(2.9), box_w - Inches(0.3), Inches(0.5),
                 "SHADE + NINJA", font_size=22, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.2), Inches(3.5), box_w - Inches(0.4), Inches(2),
                    ["Base call",
                     "Pattern-match quarters",
                     "Two-high stays two-high",
                     "Read and react"],
                    font_size=15, line_spacing=1.4)

    # Box 2
    x = start_x + box_w + gap_x
    add_shape(slide, x, Inches(2.8), box_w, box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, x + Inches(0.15), Inches(2.9), box_w - Inches(0.3), Inches(0.5),
                 "SHADE + ANGLE + Mike + NINJA", font_size=20, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.2), Inches(3.5), box_w - Inches(0.4), Inches(2),
                    ["Same pre-snap shell",
                     "DL slants to boundary",
                     "M blitzes his gap",
                     "Added pressure, same look"],
                    font_size=15, line_spacing=1.4)

    # Box 3
    x = start_x + 2 * (box_w + gap_x)
    add_shape(slide, x, Inches(2.8), box_w, box_h, FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, x + Inches(0.15), Inches(2.9), box_w - Inches(0.3), Inches(0.5),
                 "SHADE + ZEUS", font_size=22, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.2), Inches(3.5), box_w - Inches(0.4), Inches(2),
                    ["Same pre-snap shell",
                     "Zero coverage (no deep help)",
                     "Run-first, then rush",
                     "Show it, then call it off"],
                    font_size=15, line_spacing=1.4)

    # Takeaway
    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6),
                 "The offense can\u2019t diagnose the call pre-snap. That\u2019s the advantage.",
                 font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 22 — System Volume ──────────────────────────────────────────
    slide = content_slide(prs, "SYSTEM AT A GLANCE")

    data = [
        ["Category", "Count", "Detail"],
        ["4-Down Fronts", "8", "SHADE, UNDER, EYES, WIDE, DEUCES, GRIZZLY, BOSS, BOSS UNDER"],
        ["3-Down Packages", "4", "MINT, ACE, JET, SLIP"],
        ["Stunts", "11", "SLANT, ANGLE, PINCH, JACKS, SPLIT, CRASH, VEER, AA, EA, AR, ER"],
        ["Single Pressures", "4", "Mike, Will, Bandit, Dawg"],
        ["Combo Pressures", "5", "sWarM, BooM, BoW, MaD, Eat"],
        ["Packaged Pressures", "4", "Hammer, Shave, Bandit Raven, Will Raven"],
        ["Coverage Families", "4", "NINJA, Cover 1, Cover 0/Z, VIKING"],
        ["Z-Family Calls", "5", "Zeus, Zorro, Zunnel, Zill, Zike"],
        ["Alerts", "4", "BUMP, BANJO, EXCHANGE, UNDER"],
    ]
    add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.2),
              rows=10, cols=3, data=data,
              col_widths=[Inches(3.2), Inches(1.2), Inches(7.7)], font_size=16)

    add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
                 "Multiple looks. Simple, consistent rules. Layered install over 10 camp days.",
                 font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: WHAT COMES NEXT
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "WHAT COMES NEXT", "")

    # ── Slide 24 — Next Steps ─────────────────────────────────────────────
    slide = content_slide(prs, "NEXT STEPS")

    items = [
        ("Playbook ", "\u2014 Still being finished (may be more than 27 sections). Will be shared via Google Drive."),
        ("Unit meetings ", "\u2014 DL, LBs, DBs will each get deep-dives on their responsibilities."),
        ("Camp install plan ", "\u2014 10-day layered install. Probably next meeting."),
        ("Diagram documents ", "\u2014 Visual references for every front, coverage, and pressure (corrections in progress)."),
        ("Quick-reference cards ", "\u2014 M/W gap fits, NINJA teaching summary, coverage families, Bandit assignment."),
        ("This presentation ", "\u2014 Will be shared on the Drive after the meeting."),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5),
                    items, font_size=20, line_spacing=1.5)

    # ── Slide 25 — Closing ────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(2.5))
    add_gold_bar(slide, Inches(5.0))

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1),
                 "FAST.  PHYSICAL.  DISCIPLINED.", font_size=48, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.8),
                 "River Valley Vikings | 2026", font_size=24, color=GOLD,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.6),
                 "Questions?", font_size=28, color=COLUMBIA_BLUE,
                 bold=False, alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "Meeting_01_Introduction_to_the_Defense.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
