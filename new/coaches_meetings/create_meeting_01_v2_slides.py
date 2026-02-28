#!/usr/bin/env python3
"""
Generate Meeting 1 V2 slide deck: Introduction to the Defense
River Valley Vikings | 2026 Season

Changes from V1:
- New pressure-first philosophy statement (raw DC voice)
- STORM and HAVOC acronym options slide (staff picks)
- Fixed NINJA terminology (MOD/CLAMP/POACH — removed Flip/SKY/Cloud/Read 2)
- Fixed Zeus description (run-first read, not "delayed pressure")
- Added CAMP rule (whoever has contain has the QB)
- Added missing Non-Negotiable lines (#1 busts, #7 call off the horses)
- Added BOSS (F)/(B) notation
- Added "safeties echo and relay" to communication
- Added 10-day install overview slide
- Added matplotlib formation diagrams (SHADE+NINJA, SHADE+ZEUS, pressure look)
- Added 3-2-6 shell teaser (Coach Arndt collaboration)
- Section pause/discussion markers

Output: .pptx that can be uploaded directly to Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import os
import io

# ── Colors ──────────────────────────────────────────────────────────────────
COLUMBIA_BLUE = RGBColor(0x6C, 0xAC, 0xE4)
GOLD = RGBColor(0xCF, 0xA7, 0x00)
NAVY = RGBColor(0x00, 0x2D, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
FAINT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)

# Matplotlib colors (hex)
MPL_COLUMBIA = '#6CACE4'
MPL_GOLD = '#CFA700'
MPL_NAVY = '#002D62'
MPL_WHITE = '#FFFFFF'
MPL_RED = '#CC0000'
MPL_LIGHT_BG = '#F0F4F8'

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=18,
                       color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                       font_name="Calibri", line_spacing=1.3):
    """Add text with explicit line breaks as separate paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
        p.alignment = alignment
        if isinstance(line, tuple):
            # (bold_text, normal_text)
            run_b = p.add_run()
            run_b.text = line[0]
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = color
            run_b.font.bold = True
            run_b.font.name = font_name
            run_n = p.add_run()
            run_n.text = line[1]
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
            run_n.font.bold = False
            run_n.font.name = font_name
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font_name
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, bold_prefix=True, line_spacing=1.3,
                    font_name="Calibri"):
    """Add a bulleted list. Items can be strings or (bold_part, rest) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)

        if isinstance(item, tuple):
            run_b = p.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = color
            run_b.font.bold = True
            run_b.font.name = font_name

            run_n = p.add_run()
            run_n.text = item[1]
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
            run_n.font.bold = False
            run_n.font.name = font_name
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = False
            run.font.name = font_name

    return txBox


def add_gold_bar(slide, top, width=SLIDE_W):
    """Thin gold accent bar."""
    add_shape(slide, Inches(0), top, width, Pt(4), GOLD)


def section_header_slide(prs, title, subtitle=None):
    """Navy background section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(2.8))
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
                 title, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1),
                     subtitle, font_size=22, color=COLUMBIA_BLUE, bold=False,
                     alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title):
    """White background content slide with navy top bar and gold accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_gold_bar(slide, Inches(1.1))
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.85),
                 title, font_size=32, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(5), Inches(0.4),
                 "River Valley Vikings | 2026 Defense", font_size=10,
                 color=MED_GRAY, bold=False)
    return slide


def discussion_slide(prs, title="DISCUSSION", prompt="Questions before we move on?"):
    """Brief pause slide for staff discussion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(3.0))
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
                 title, font_size=40, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(1),
                 prompt, font_size=24, color=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)
    return slide


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None, header_color=NAVY, header_text_color=WHITE,
              font_size=14, row_alt_color=FAINT_BLUE):
    """Add a formatted table. data is list of lists (rows x cols)."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ── Diagram Helpers ─────────────────────────────────────────────────────────

def _setup_field_ax(ax, title="", width=53.3, depth=25):
    """Configure a clean football-field-style axes."""
    ax.set_xlim(-width/2 - 2, width/2 + 2)
    ax.set_ylim(-3, depth + 2)
    ax.set_aspect('equal')
    ax.set_facecolor(MPL_LIGHT_BG)
    ax.axhline(y=0, color='#AAAAAA', linewidth=1.5, linestyle='-')  # LOS
    ax.axhline(y=-0.3, color=MPL_NAVY, linewidth=3)  # Thick LOS marker
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    if title:
        ax.set_title(title, fontsize=16, fontweight='bold', color=MPL_NAVY,
                      pad=12)


def _draw_player(ax, x, y, label, color=MPL_COLUMBIA, text_color=MPL_WHITE,
                 shape='circle', size=1.4):
    """Draw a player marker with label."""
    if shape == 'circle':
        circle = plt.Circle((x, y), size, facecolor=color, edgecolor=MPL_NAVY,
                             linewidth=1.5, zorder=5)
        ax.add_patch(circle)
    elif shape == 'square':
        sq = mpatches.FancyBboxPatch((x - size, y - size), size*2, size*2,
                                      boxstyle="round,pad=0.15",
                                      facecolor=color, edgecolor=MPL_NAVY,
                                      linewidth=1.5, zorder=5)
        ax.add_patch(sq)
    elif shape == 'diamond':
        diamond = plt.Polygon([(x, y+size*1.2), (x+size, y),
                                (x, y-size*1.2), (x-size, y)],
                               facecolor=color, edgecolor=MPL_NAVY,
                               linewidth=1.5, zorder=5)
        ax.add_patch(diamond)

    ax.text(x, y, label, ha='center', va='center', fontsize=9,
            fontweight='bold', color=text_color, zorder=6)


def _draw_ol_with_shading(ax, front='SHADE'):
    """Draw OL with shading to indicate DL technique alignment.
    Adds small triangles/markers on OL to show where DL aligns."""
    ol_color = '#888888'
    ol_labels = ['LT', 'LG', 'C', 'RG', 'RT']
    ol_xs = [-8, -4, 0, 4, 8]
    # Map fronts to technique shading positions (field=right, boundary=left)
    # Each entry: (ol_index, side) where side='inside','outside','head','shade'
    shade_map = {
        'SHADE': [
            # A=5 on RT(field): outside shade RT
            (4, 'outside', 'A'),
            # T=3 on RG(field): outside shade RG
            (3, 'outside', 'T'),
            # N=2i on LG(boundary): inside shade LG
            (2, 'inside_left', 'N'),
            # E=5 on LT(boundary): outside shade LT (outside = further left)
            (0, 'outside_left', 'E'),
        ],
    }

    for x_pos, lbl in zip(ol_xs, ol_labels):
        sq = mpatches.FancyBboxPatch((x_pos - 1.1, -2.5), 2.2, 2.2,
                                      boxstyle="round,pad=0.1",
                                      facecolor=ol_color, edgecolor='#555555',
                                      linewidth=1, zorder=3)
        ax.add_patch(sq)
        ax.text(x_pos, -1.4, lbl, ha='center', va='center', fontsize=7,
                color=MPL_WHITE, fontweight='bold', zorder=4)

    # Add technique shade markers on the OL
    if front in shade_map:
        for ol_idx, side, dl_lbl in shade_map[front]:
            ox = ol_xs[ol_idx]
            if side == 'outside':
                # Right side of the OL (field side)
                marker_x = ox + 1.3
            elif side == 'outside_left':
                # Left side of the OL (boundary side, outside = more left)
                marker_x = ox - 1.3
            elif side == 'inside_left':
                # Inside shade on left side
                marker_x = ox + 0.8
            elif side == 'inside':
                marker_x = ox - 0.8
            else:
                marker_x = ox
            # Small triangle marker showing technique shade
            tri = plt.Polygon([(marker_x, -0.1), (marker_x - 0.4, 0.5),
                                (marker_x + 0.4, 0.5)],
                               facecolor=MPL_COLUMBIA, edgecolor=MPL_NAVY,
                               linewidth=0.8, zorder=4, alpha=0.7)
            ax.add_patch(tri)


def _draw_offense(ax, has_te=False, trips=False, front='SHADE'):
    """Draw a standard offensive formation with OL shading for DL techniques."""
    ol_color = '#888888'

    # Draw OL with technique shading
    _draw_ol_with_shading(ax, front)

    # QB — darker color for visibility
    qb_color = '#444444'
    sq = mpatches.FancyBboxPatch((-1.3, -5.7), 2.6, 2.4,
                                  boxstyle="round,pad=0.1",
                                  facecolor=qb_color, edgecolor='#222222',
                                  linewidth=1.5, zorder=3)
    ax.add_patch(sq)
    ax.text(0, -4.5, 'QB', ha='center', va='center', fontsize=8,
            color=MPL_WHITE, fontweight='bold', zorder=4)

    if not trips:
        # 2x2 WRs
        for wx in [-24, -14, 14, 24]:
            sq = mpatches.FancyBboxPatch((wx - 1, -2.3), 2, 2,
                                          boxstyle="round,pad=0.1",
                                          facecolor=ol_color, edgecolor='#555555',
                                          linewidth=1, zorder=3)
            ax.add_patch(sq)
            lbl = 'WR' if abs(wx) > 20 else 'SL'
            ax.text(wx, -1.3, lbl, ha='center', va='center', fontsize=6,
                    color=MPL_WHITE, fontweight='bold', zorder=4)
        # RB — darker color for visibility
        sq = mpatches.FancyBboxPatch((-1.3, -8.7), 2.6, 2.4,
                                      boxstyle="round,pad=0.1",
                                      facecolor=qb_color, edgecolor='#222222',
                                      linewidth=1.5, zorder=3)
        ax.add_patch(sq)
        ax.text(0, -7.5, 'RB', ha='center', va='center', fontsize=8,
                color=MPL_WHITE, fontweight='bold', zorder=4)
    else:
        # 3x1 trips to field (right side)
        for wx in [-24]:
            sq = mpatches.FancyBboxPatch((wx - 1, -2.3), 2, 2,
                                          boxstyle="round,pad=0.1",
                                          facecolor=ol_color, edgecolor='#555555',
                                          linewidth=1, zorder=3)
            ax.add_patch(sq)
            ax.text(wx, -1.3, 'WR', ha='center', va='center', fontsize=6,
                    color=MPL_WHITE, fontweight='bold', zorder=4)
        trips_xs = [14, 17, 24]
        trips_lbls = ['SL', 'SL', 'WR']
        for wx, lbl in zip(trips_xs, trips_lbls):
            sq = mpatches.FancyBboxPatch((wx - 1, -2.3), 2, 2,
                                          boxstyle="round,pad=0.1",
                                          facecolor=ol_color, edgecolor='#555555',
                                          linewidth=1, zorder=3)
            ax.add_patch(sq)
            ax.text(wx, -1.3, lbl, ha='center', va='center', fontsize=6,
                    color=MPL_WHITE, fontweight='bold', zorder=4)
        # RB — darker
        sq = mpatches.FancyBboxPatch((-1.3, -8.7), 2.6, 2.4,
                                      boxstyle="round,pad=0.1",
                                      facecolor=qb_color, edgecolor='#222222',
                                      linewidth=1.5, zorder=3)
        ax.add_patch(sq)
        ax.text(0, -7.5, 'RB', ha='center', va='center', fontsize=8,
                color=MPL_WHITE, fontweight='bold', zorder=4)


def _add_zone_label(ax, x, y, text, color=MPL_NAVY, fontsize=8):
    """Add a small zone/assignment label."""
    ax.text(x, y, text, ha='center', va='center', fontsize=fontsize,
            color=color, fontstyle='italic', zorder=7)


def create_shade_ninja_diagram():
    """SHADE + NINJA vs 2x2 — base look."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6.5))
    _setup_field_ax(ax, "SHADE + NINJA  (Base)")
    _draw_offense(ax)

    # -- Defense --
    # DL: A(5 field), T(3 field), N(2i bnd), E(5 bnd)
    # Field = right, Boundary = left
    _draw_player(ax, 11, 1.5, 'A', MPL_COLUMBIA)    # 5-tech field (outside OT)
    _draw_player(ax, 5.5, 1.5, 'T', MPL_COLUMBIA)   # 3-tech field (outside OG)
    _draw_player(ax, -3, 1.5, 'N', MPL_COLUMBIA)     # 2i bnd (inside OG)
    _draw_player(ax, -11, 1.5, 'E', MPL_COLUMBIA)    # 5-tech bnd

    # LBs: M(field A gap area), W(boundary B gap area)
    _draw_player(ax, 3, 5, 'M', MPL_GOLD, MPL_NAVY)
    _draw_player(ax, -6, 5, 'W', MPL_GOLD, MPL_NAVY)

    # B (Bandit) — apex field
    _draw_player(ax, 14, 5, 'B', MPL_GOLD, MPL_NAVY)

    # Safeties — two-high
    _draw_player(ax, 10, 13, 'FS', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -10, 13, 'D', MPL_NAVY, MPL_WHITE)

    # Corners
    _draw_player(ax, 24, 7, 'FC', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -24, 7, 'BC', MPL_NAVY, MPL_WHITE)

    # Zone labels
    _add_zone_label(ax, 24, 10, 'MOD: #1')
    _add_zone_label(ax, 10, 16, 'MOD: #2\u2192#1')
    _add_zone_label(ax, 14, 8, 'Flat/#2')
    _add_zone_label(ax, -24, 10, 'CLAMP: keys #2')
    _add_zone_label(ax, -10, 16, 'CLAMP: #2\u2192#1')
    _add_zone_label(ax, 3, 8, 'Hook')
    _add_zone_label(ax, -6, 8, 'Hook/#3')

    # Field / Boundary labels
    ax.text(26, 20, 'FIELD \u2192', ha='right', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')
    ax.text(-26, 20, '\u2190 BOUNDARY', ha='left', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')

    # Two-high shell bracket
    ax.annotate('', xy=(13, 14.5), xytext=(-13, 14.5),
                arrowprops=dict(arrowstyle='<->', color=MPL_GOLD, lw=1.5))
    ax.text(0, 15.5, 'TWO-HIGH SHELL', ha='center', va='center',
            fontsize=10, color=MPL_GOLD, fontweight='bold')

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor=MPL_LIGHT_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def create_shade_zeus_diagram():
    """SHADE + ZEUS vs 2x2 — run-first zero look."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6.5))
    _setup_field_ax(ax, "SHADE + ZEUS  (Run-First Zero)")
    _draw_offense(ax)

    # DL — same as SHADE
    _draw_player(ax, 11, 1.5, 'A', MPL_COLUMBIA)
    _draw_player(ax, 5.5, 1.5, 'T', MPL_COLUMBIA)
    _draw_player(ax, -3, 1.5, 'N', MPL_COLUMBIA)
    _draw_player(ax, -11, 1.5, 'E', MPL_COLUMBIA)

    # LBs — labels ABOVE the circles to avoid overlap
    _draw_player(ax, 3, 5, 'M', MPL_GOLD, MPL_NAVY)
    _draw_player(ax, -6, 5, 'W', MPL_GOLD, MPL_NAVY)

    # B — apex, MEG on RB
    _draw_player(ax, 14, 5, 'B', MPL_GOLD, MPL_NAVY)

    # Safeties — two-high pre-snap (same shell!)
    _draw_player(ax, 10, 13, 'FS', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -10, 13, 'D', MPL_NAVY, MPL_WHITE)

    # Corners
    _draw_player(ax, 24, 7, 'FC', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -24, 7, 'BC', MPL_NAVY, MPL_WHITE)

    # Assignment labels — Zeus specific
    _add_zone_label(ax, 24, 10, 'Man #1')
    _add_zone_label(ax, -24, 10, 'Man #1')
    _add_zone_label(ax, 10, 16, 'Man #2 field')
    _add_zone_label(ax, -10, 16, 'Man #2 bnd')

    # M/B/W run-first indicators — positioned well ABOVE the player circles
    ax.text(3, 8.2, 'Run fit \u2192\nAdd-on rush', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold', zorder=10)
    ax.text(14, 8.2, 'Run fit \u2192\nMEG on RB', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold', zorder=10)
    ax.text(-6, 8.2, 'Run fit \u2192\nRB funnel', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold', zorder=10)

    # Two-high shell label
    ax.annotate('', xy=(13, 14.5), xytext=(-13, 14.5),
                arrowprops=dict(arrowstyle='<->', color=MPL_GOLD, lw=1.5))
    ax.text(0, 15.5, 'SAME TWO-HIGH SHELL PRE-SNAP', ha='center',
            va='center', fontsize=10, color=MPL_GOLD, fontweight='bold')

    ax.text(0, 19, 'ZERO COVERAGE \u2014 NO DEEP HELP\nRun first. Pass read \u2192 rush / take RB.',
            ha='center', va='center', fontsize=11, color=MPL_RED,
            fontweight='bold')

    # Field / Boundary labels
    ax.text(26, 20, 'FIELD \u2192', ha='right', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')
    ax.text(-26, 20, '\u2190 BOUNDARY', ha='left', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor=MPL_LIGHT_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def create_shade_oregon_diagram():
    """SHADE + OREGON — Cover 1 with FS as post safety."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6.5))
    _setup_field_ax(ax, "SHADE + OREGON  (Cover 1 \u2014 FS = Post)")
    _draw_offense(ax)

    # DL — same as SHADE
    _draw_player(ax, 11, 1.5, 'A', MPL_COLUMBIA)
    _draw_player(ax, 5.5, 1.5, 'T', MPL_COLUMBIA)
    _draw_player(ax, -3, 1.5, 'N', MPL_COLUMBIA)
    _draw_player(ax, -11, 1.5, 'E', MPL_COLUMBIA)

    # LBs — M and W funnel RB
    _draw_player(ax, 3, 5, 'M', MPL_GOLD, MPL_NAVY)
    _draw_player(ax, -6, 5, 'W', MPL_GOLD, MPL_NAVY)

    # B — apex field, man #2 field
    _draw_player(ax, 14, 5, 'B', MPL_GOLD, MPL_NAVY)

    # FS — post safety, deep middle (rotates from two-high pre-snap)
    _draw_player(ax, 0, 15, 'FS', MPL_NAVY, MPL_WHITE, size=1.6)

    # D — man #2 boundary
    _draw_player(ax, -12, 8, 'D', MPL_NAVY, MPL_WHITE)

    # Corners
    _draw_player(ax, 24, 7, 'FC', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -24, 7, 'BC', MPL_NAVY, MPL_WHITE)

    # Assignment labels
    _add_zone_label(ax, 24, 10, 'Man #1 field')
    _add_zone_label(ax, -24, 10, 'Man #1 bnd')
    _add_zone_label(ax, 14, 8, 'Man #2 field')
    _add_zone_label(ax, -12, 11, 'Man #2 bnd')

    # FS post label
    ax.text(4, 16, 'POST / MOF\n(Free Safety)', fontsize=9,
            color=MPL_RED, ha='left', fontweight='bold')

    # M/W labels — RB funnel
    ax.text(3, 7.8, 'RB funnel\n(to you = take)', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold')
    ax.text(-6, 7.8, 'RB funnel\n(away = rush)', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold')

    # B label
    ax.text(18.5, 5.5, 'Man #2 field\nor #3 away', fontsize=8,
            color=MPL_RED, ha='left', fontweight='bold')

    ax.text(0, 19, 'ONE FREE SAFETY \u2014 MAN COVERAGE + POST HELP',
            ha='center', va='center', fontsize=11, color=MPL_RED,
            fontweight='bold')

    # Rotation arrow showing FS moving to MOF
    ax.annotate('', xy=(0, 13.5), xytext=(8, 12),
                arrowprops=dict(arrowstyle='->', color=MPL_GOLD, lw=2,
                                linestyle='dashed'))
    ax.text(9, 12, 'Pre-snap\ntwo-high', fontsize=7,
            color=MPL_GOLD, ha='left', fontweight='bold')

    ax.text(26, 20, 'FIELD \u2192', ha='right', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')
    ax.text(-26, 20, '\u2190 BOUNDARY', ha='left', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor=MPL_LIGHT_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def create_pressure_look_diagram():
    """SHADE + ANGLE + Mike + ZEUS — pressure from disguise."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6.5))
    _setup_field_ax(ax, "SHADE + ANGLE + Mike + ZEUS")
    _draw_offense(ax)

    # DL — SHADE base alignment, ANGLE = slant to boundary
    _draw_player(ax, 11, 1.5, 'A', MPL_COLUMBIA)
    _draw_player(ax, 5.5, 1.5, 'T', MPL_COLUMBIA)
    _draw_player(ax, -3, 1.5, 'N', MPL_COLUMBIA)
    _draw_player(ax, -11, 1.5, 'E', MPL_COLUMBIA)

    # ANGLE stunt: T+N slant to boundary (left)
    # T slants from 3-tech (outside RG, x=5.5) through B gap toward A gap area
    # Arrow angled at the gap between RG(x=4) and C(x=0) — target A gap (x=2)
    ax.annotate('', xy=(2, -0.3), xytext=(5.5, 1.5),
                arrowprops=dict(arrowstyle='->', color=MPL_RED, lw=2.5))
    # N slants from 2i (inside LG, x=-3) through A gap toward B gap boundary
    # Arrow angled at the gap between LG(x=-4) and LT(x=-8) — target B gap (x=-6)
    ax.annotate('', xy=(-6, -0.3), xytext=(-3, 1.5),
                arrowprops=dict(arrowstyle='->', color=MPL_RED, lw=2.5))

    # M blitz arrow — Mike blitzes the B GAP field (between RG and RT)
    # B gap field is between RG(x=4) and RT(x=8), so target ~x=6
    _draw_player(ax, 3, 5, 'M', MPL_GOLD, MPL_NAVY)
    ax.annotate('', xy=(6, -0.3), xytext=(3, 3.5),
                arrowprops=dict(arrowstyle='->', color=MPL_RED, lw=2.5))
    ax.text(7, 4, 'MIKE BLITZ\n(B gap field)', fontsize=9, color=MPL_RED,
            fontweight='bold')

    # W — run fit / RB funnel (Zeus coverage, not dropping into zone)
    _draw_player(ax, -6, 5, 'W', MPL_GOLD, MPL_NAVY)
    ax.text(-6, 7.5, 'Run fit \u2192\nRB funnel', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold')

    # B — MEG on RB (Zeus coverage)
    _draw_player(ax, 14, 5, 'B', MPL_GOLD, MPL_NAVY)
    ax.text(14, 7.5, 'Run fit \u2192\nMEG on RB', fontsize=8,
            color=MPL_RED, ha='center', fontweight='bold')

    # Safeties — two-high
    _draw_player(ax, 10, 13, 'FS', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -10, 13, 'D', MPL_NAVY, MPL_WHITE)

    # Corners
    _draw_player(ax, 24, 7, 'FC', MPL_NAVY, MPL_WHITE)
    _draw_player(ax, -24, 7, 'BC', MPL_NAVY, MPL_WHITE)

    # Assignment labels — Zeus
    _add_zone_label(ax, 24, 10, 'Man #1')
    _add_zone_label(ax, -24, 10, 'Man #1')
    _add_zone_label(ax, 10, 16, 'Man #2 field')
    _add_zone_label(ax, -10, 16, 'Man #2 bnd')

    # Labels
    _add_zone_label(ax, 0, 19, 'SAME PRE-SNAP SHELL \u2014 DL SLANTS + MIKE BLITZES POST-SNAP')

    ax.text(26, 20, 'FIELD \u2192', ha='right', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')
    ax.text(-26, 20, '\u2190 BOUNDARY', ha='left', va='center', fontsize=11,
            color=MPL_GOLD, fontweight='bold')

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor=MPL_LIGHT_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


def create_three_look_comparison():
    """Side-by-side: SHADE+NINJA, SHADE+ZEUS, SHADE+OREGON — shows shell variety."""
    fig, axes = plt.subplots(1, 3, figsize=(16, 5.5))

    titles = ["SHADE + NINJA", "SHADE + ZEUS", "SHADE + OREGON"]
    subtitles = ["Two-High / Quarters", "Two-High \u2192 Zero", "One-High / Man + Free"]

    for idx, (ax, title, sub) in enumerate(zip(axes, titles, subtitles)):
        _setup_field_ax(ax, "", depth=18)

        # Simplified offense — OL with technique shading
        ol_xs = [-6, -3, 0, 3, 6]
        ol_labels = ['LT', 'LG', 'C', 'RG', 'RT']
        for x_pos, ol_lbl in zip(ol_xs, ol_labels):
            sq = mpatches.FancyBboxPatch((x_pos - 0.8, -2), 1.6, 1.6,
                                          boxstyle="round,pad=0.1",
                                          facecolor='#888888',
                                          edgecolor='#555555',
                                          linewidth=1, zorder=3)
            ax.add_patch(sq)
            ax.text(x_pos, -1.2, ol_lbl, ha='center', va='center', fontsize=5,
                    color=MPL_WHITE, fontweight='bold', zorder=4)
        # SHADE technique shading markers on OL
        # A=5 on RT (outside field), T=3 on RG (outside field)
        # N=2i on LG (inside toward C), E=5 on LT (outside boundary)
        shade_markers = [(6, 'right'), (3, 'right'), (0, 'left_inside'), (-6, 'left')]
        for mx, side in shade_markers:
            if side == 'right':
                marker_x = mx + 1.0
            elif side == 'left':
                marker_x = mx - 1.0
            elif side == 'left_inside':
                marker_x = mx + 0.6
            else:
                marker_x = mx
            tri = plt.Polygon([(marker_x, -0.2), (marker_x - 0.3, 0.3),
                                (marker_x + 0.3, 0.3)],
                               facecolor=MPL_COLUMBIA, edgecolor=MPL_NAVY,
                               linewidth=0.6, zorder=4, alpha=0.7)
            ax.add_patch(tri)

        # QB + RB — high contrast for visibility
        qb_c = '#333333'
        sq = mpatches.FancyBboxPatch((-0.9, -4.7), 1.8, 1.8,
                                      boxstyle="round,pad=0.1",
                                      facecolor=qb_c, edgecolor='#111111',
                                      linewidth=1.5, zorder=3)
        ax.add_patch(sq)
        ax.text(0, -3.8, 'QB', ha='center', va='center', fontsize=7,
                color=MPL_WHITE, fontweight='bold', zorder=4)

        sq2 = mpatches.FancyBboxPatch((-0.9, -7.2), 1.8, 1.8,
                                       boxstyle="round,pad=0.1",
                                       facecolor=qb_c, edgecolor='#111111',
                                       linewidth=1.5, zorder=3)
        ax.add_patch(sq2)
        ax.text(0, -6.3, 'RB', ha='center', va='center', fontsize=7,
                color=MPL_WHITE, fontweight='bold', zorder=4)

        # WRs
        for wx in [-18, -10, 10, 18]:
            sq = mpatches.FancyBboxPatch((wx - 0.7, -1.8), 1.4, 1.4,
                                          boxstyle="round,pad=0.1",
                                          facecolor='#999999',
                                          edgecolor='#666666',
                                          linewidth=0.8, zorder=3)
            ax.add_patch(sq)

        # DL — same for all
        for dx, dl in [(8, 'A'), (4, 'T'), (-2.5, 'N'), (-8, 'E')]:
            _draw_player(ax, dx, 1.2, dl, MPL_COLUMBIA, size=1.1)

        # Corners — same for all
        _draw_player(ax, 18, 5, 'FC', MPL_NAVY, MPL_WHITE, size=1.1)
        _draw_player(ax, -18, 5, 'BC', MPL_NAVY, MPL_WHITE, size=1.1)

        if idx == 0:  # NINJA — two-high stays two-high
            _draw_player(ax, 7, 11, 'FS', MPL_NAVY, MPL_WHITE, size=1.1)
            _draw_player(ax, -7, 11, 'D', MPL_NAVY, MPL_WHITE, size=1.1)
            _draw_player(ax, 2, 4, 'M', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, -4, 4, 'W', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, 10, 4, 'B', MPL_GOLD, MPL_NAVY, size=1.1)
        elif idx == 1:  # ZEUS — two-high pre-snap, zero post-snap
            _draw_player(ax, 7, 11, 'FS', MPL_NAVY, MPL_WHITE, size=1.1)
            _draw_player(ax, -7, 11, 'D', MPL_NAVY, MPL_WHITE, size=1.1)
            _draw_player(ax, 2, 4, 'M', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, -4, 4, 'W', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, 10, 4, 'B', MPL_GOLD, MPL_NAVY, size=1.1)
        else:  # OREGON — FS rotates to MOF post-snap (one-high)
            # FS stays deep middle (post-snap MOF)
            _draw_player(ax, 0, 13, 'FS', MPL_NAVY, MPL_WHITE, size=1.1)
            # D drops down to man #2 boundary
            _draw_player(ax, -10, 6, 'D', MPL_NAVY, MPL_WHITE, size=1.1)
            # B at apex, man #2 field
            _draw_player(ax, 10, 4, 'B', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, 2, 4, 'M', MPL_GOLD, MPL_NAVY, size=1.1)
            _draw_player(ax, -4, 4, 'W', MPL_GOLD, MPL_NAVY, size=1.1)
            # Show rotation arrow
            ax.annotate('POST', xy=(0, 14.5), fontsize=7, color=MPL_RED,
                        ha='center', fontweight='bold')

        # Title and subtitle below diagram
        ax.text(0, 17, title, ha='center', va='center', fontsize=12,
                color=MPL_NAVY, fontweight='bold')
        ax.text(0, 15.5, sub, ha='center', va='center', fontsize=9,
                color=MPL_RED, fontstyle='italic')

    fig.suptitle("THREE SHELLS  \u2014  THREE DIFFERENT ANSWERS",
                 fontsize=14, fontweight='bold', color=MPL_NAVY, y=0.98)
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor=MPL_LIGHT_BG)
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Build Presentation ──────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Pre-generate diagrams
    ninja_img = create_shade_ninja_diagram()
    zeus_img = create_shade_zeus_diagram()
    oregon_img = create_shade_oregon_diagram()
    pressure_img = create_pressure_look_diagram()
    comparison_img = create_three_look_comparison()

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_gold_bar(slide, Inches(1.8))
    add_gold_bar(slide, Inches(5.3))

    add_text_box(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(1.2),
                 "RIVER VALLEY VIKINGS", font_size=52, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.8),
                 "2026 DEFENSIVE SYSTEM", font_size=36, color=COLUMBIA_BLUE,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.8),
                 "Coaches Meeting 1  |  Introduction to the Defense",
                 font_size=22, color=GOLD, bold=False,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.7),
                 "Confidential \u2014 Coaching Staff Only", font_size=16,
                 color=MED_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Agenda
    # ════════════════════════════════════════════════════════════════════════
    slide = content_slide(prs, "TODAY\u2019S AGENDA")

    items = [
        ("1.  ", "Defensive Identity \u2014 Who We Are"),
        ("2.  ", "Position Designations \u2014 Our Language"),
        ("3.  ", "Technique & Gap Names \u2014 Alignment Language"),
        ("4.  ", "Call Grammar \u2014 How We Communicate"),
        ("5.  ", "The System Overview \u2014 What\u2019s in the Toolbox"),
        ("6.  ", "How It Fits Together \u2014 Diagrams"),
        ("7.  ", "Next Steps"),
    ]
    add_bullet_list(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(5),
                    items, font_size=24, line_spacing=1.5)

    add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6),
                 "Goal: Walk out speaking the same language, knowing the system, and ready to teach it.",
                 font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.LEFT)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: DEFENSIVE IDENTITY
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "DEFENSIVE IDENTITY", "Who We Are")

    # ── Philosophy ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(1.2))
    add_gold_bar(slide, Inches(6.3))

    add_text_box(slide, Inches(1.2), Inches(0.3), Inches(10.9), Inches(0.8),
                 "OUR PHILOSOPHY", font_size=28, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    philosophy_lines = [
        "We\u2019re a pressure defense. We attack with plus-one advantages",
        "and play zero behind it. When we\u2019re not pressuring,",
        "we play man. If they pick it up, we change the picture",
        "and come again.",
        "",
        "We don\u2019t sit back and wait for offenses to beat us. We attack.",
        "We set edges, we own our gaps, and we get 11 hats to the ball.",
        "",
        "No explosives. No free yards. Make them earn every first down,",
        "and when they don\u2019t, take the ball.",
    ]
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.6),
                                      Inches(10.9), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(philosophy_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(36)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(24)
        run.font.color.rgb = WHITE
        run.font.italic = True
        run.font.name = "Calibri"

    # ── Base Structure ───────────────────────────────────────────────────
    slide = content_slide(prs, "BASE STRUCTURE")

    data = [
        ["Element", "Detail"],
        ["Personnel", "4-2-5 base; 3-4 / 3-2-6 packages"],
        ["Strength", "Field"],
        ["Shell", "Two-high"],
        ["Box", "6-man"],
        ["Call Format", "Front + Stunt + Blitz + Coverage/Tags"],
        ["Communication", "Signaled front/stunt/blitz; coverage is verbal. Safeties echo and relay."],
    ]
    add_table(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5),
              rows=7, cols=2, data=data,
              col_widths=[Inches(3), Inches(7.9)], font_size=18)

    add_text_box(slide, Inches(1.2), Inches(6.3), Inches(10.9), Inches(0.5),
                 "Multiple looks, simple rules. The picture changes \u2014 the principles don\u2019t.",
                 font_size=16, color=NAVY, bold=True)

    # ── Non-Negotiables ──────────────────────────────────────────────────
    slide = content_slide(prs, "THE NON-NEGOTIABLES")

    items = [
        ("1. No explosives. ", "Make them snap it again. Eliminate busts and cheap yards."),
        ("2. Aggressive with discipline. ", "Attack without losing leverage or committing penalties."),
        ("3. Box the run. ", "Set hard edges, compress gaps, keep the ball in help."),
        ("4. Relentless pursuit. ", "11 to the ball, correct angles, swarming tackles."),
        ("5. Win situations. ", "1st down sets the series; 3rd down = get off the field."),
        ("6. Smart aggression. ", "Plus-one rushers with zero behind it. No-win for the offense. Call off the dawgs when needed."),
        ("7. Communication standard. ", "Simple, loud, early. Everyone echoes. Alerts: BUMP / BANJO / EXCHANGE / UNDER."),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5),
                    items, font_size=18, line_spacing=1.35)

    # ── The Standard ─────────────────────────────────────────────────────
    slide = content_slide(prs, "THE STANDARD")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
                 "The baseline. Every snap, every rep, every day.",
                 font_size=18, color=MED_GRAY)

    items = [
        ("No busted coverages ", "\u2014 know your call, know your rules, communicate"),
        ("Stay in phase ", "\u2014 maintain leverage in man; no trailing, no freelancing"),
        ("No free routes ", "\u2014 collision every release; reroute, disrupt timing, deny easy access"),
        ("No lost edges ", "\u2014 contain player keeps QB and ball inside"),
        ("No uncovered gaps ", "\u2014 adjust instantly if gap changes with stunt/pressure"),
        ("No missed tackles ", "\u2014 leverage + near foot + wrap; eliminate YAC"),
        ("No penalties ", "that extend drives"),
        ("No loafs ", "\u2014 pursuit is mandatory"),
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.5),
                    items, font_size=20, line_spacing=1.4)

    # ── STORM vs HAVOC ───────────────────────────────────────────────────
    slide = content_slide(prs, "WHAT WE WILL LIVE BY")

    # STORM box
    box = add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.0),
                    FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.7),
                 "STORM", font_size=36, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.4), Inches(0.5),
                 '"Bring the STORM."', font_size=18, color=MED_GRAY,
                 bold=False, alignment=PP_ALIGN.CENTER)

    storm_items = [
        ("S ", "\u2014 Sound. Gap integrity, assignment football."),
        ("T ", "\u2014 Tough. Physical at the point, finish every rep."),
        ("O ", "\u2014 Overwhelming. Plus-one advantage, pressure from every angle."),
        ("R ", "\u2014 Relentless. 11 to the ball, pursuit is mandatory."),
        ("M ", "\u2014 Multiple. Same shell, different answers, keep them guessing."),
    ]
    add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(3.5),
                    storm_items, font_size=16, line_spacing=1.5)

    # HAVOC box
    box2 = add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.0),
                     FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.7),
                 "HAVOC", font_size=36, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5.4), Inches(0.5),
                 '"Create HAVOC."', font_size=18, color=MED_GRAY,
                 bold=False, alignment=PP_ALIGN.CENTER)

    havoc_items = [
        ("H ", "\u2014 Hit. Collision every route, attack the ball carrier."),
        ("A ", "\u2014 Attack. Pressure with purpose. Zeus, Zorro, zero."),
        ("V ", "\u2014 Violence at the point. Win at the LOS, dominate gaps."),
        ("O ", "\u2014 Overwhelm. Multiple looks, disguise, put them in conflict."),
        ("C ", "\u2014 Contain. No lost edges, no free runners, discipline in chaos."),
    ]
    add_bullet_list(slide, Inches(7.2), Inches(2.9), Inches(5.2), Inches(3.5),
                    havoc_items, font_size=16, line_spacing=1.5)

    add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
                 "This is what we stand for. Pick the one that fits.",
                 font_size=16, color=NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: POSITION DESIGNATIONS
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "POSITION DESIGNATIONS", "Our Language")

    # ── Position Table ───────────────────────────────────────────────────
    slide = content_slide(prs, "THE 11 POSITIONS")

    data = [
        ["Letter", "Name", "Attributes"],
        ["FC", "Field Corner", "Speed, cover skills, tackling in space"],
        ["FS", "Field Safety", "Range, communication, run support, ball skills"],
        ["B", "Bandit (Nickel)", "Versatile \u2014 DB in coverage, OLB in 3-down; physical, athletic"],
        ["A", "Anchor (Field DE)", "Edge setter, pass rusher, QB player (CAMP)"],
        ["T", "Tackle (Field DT)", "Power, interior push, gap control"],
        ["N", "Nose (Boundary DT)", "Anchor, interior push, gap control"],
        ["E", "Edge (Boundary DE)", "Edge setter, pass rusher; OLB in 3-down"],
        ["M", "Mike (Field ILB)", "Instincts, physicality, coverage range"],
        ["W", "Will (Boundary ILB)", "Instincts, physicality, coverage range"],
        ["D", "Dawg (Boundary S)", "Physicality, run support, communication, coverage"],
        ["BC", "Boundary Corner", "Physical, press ability, tackling"],
    ]
    add_table(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.8),
              rows=12, cols=3, data=data,
              col_widths=[Inches(1.5), Inches(3.5), Inches(7.1)], font_size=16)

    # ── Defaults & Key Roles ─────────────────────────────────────────────
    slide = content_slide(prs, "DEFAULT SIDES & KEY RULES")

    # Field side box
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2),
              FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.1), Inches(0.5),
                 "FIELD SIDE (wide side)", font_size=22, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(5.1), Inches(1.4),
                    ["A (Anchor) \u2014 Field DE",
                     "T (Tackle) \u2014 Field DT",
                     "B (Bandit) \u2014 Nickel / Field force",
                     "FC / FS \u2014 Field corner & safety"],
                    font_size=16, color=DARK_GRAY)

    # Boundary side box
    add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.2),
              FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5.1), Inches(0.5),
                 "BOUNDARY SIDE (short side)", font_size=22, color=NAVY,
                 bold=True)
    add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.1), Inches(1.4),
                    ["E (Edge) \u2014 Boundary DE",
                     "N (Nose) \u2014 Boundary DT",
                     "D (Dawg) \u2014 Boundary safety / force",
                     "BC \u2014 Boundary corner"],
                    font_size=16, color=DARK_GRAY)

    # Key distinctions
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
                 "Key Rules", font_size=22, color=NAVY, bold=True)
    items = [
        ("Contain: ", "Maintain outside leverage. Squeeze run inside. Set the edge."),
        ("Force: ", "Collision at or behind LOS. Turn the ball inside. More aggressive than contain."),
        ("CAMP: ", "A \u2018CAMP rule\u2019 is a default rule \u2014 it\u2019s always on unless told otherwise. Example: whoever has contain has the QB."),
        ("Safeties are fixers. ", "FS owns field-side calls. D owns boundary-side calls. They echo and relay."),
    ]
    add_bullet_list(slide, Inches(1.0), Inches(4.7), Inches(11.5), Inches(2.2),
                    items, font_size=16, line_spacing=1.4)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: TECHNIQUE & GAP TERMINOLOGY
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "TECHNIQUE & GAP NAMES", "Alignment Language")

    # ── DL Technique Chart ───────────────────────────────────────────────
    slide = content_slide(prs, "DL TECHNIQUE CHART")

    data = [
        ["Tech", "Alignment"],
        ["0", "Head up on Center"],
        ["1", "Shade on Center"],
        ["2i", "Inside shade on Guard"],
        ["2", "Head up on Guard"],
        ["3", "Outside shade on Guard"],
        ["4i", "Inside shade on Tackle"],
        ["4", "Head up on Tackle"],
        ["5", "Outside shade on Tackle"],
        ["6i", "Inside shade on TE"],
        ["6", "Head up on TE"],
        ["7", "Outside shade on TE"],
        ["8i", "Inside shade on Wing"],
        ["8", "Head up on Wing"],
        ["9", "Outside shade on Wing"],
    ]
    add_table(slide, Inches(1.5), Inches(1.35), Inches(10.3), Inches(5.8),
              rows=15, cols=2, data=data,
              col_widths=[Inches(2), Inches(8.3)], font_size=16)

    # ── Gap Names ────────────────────────────────────────────────────────
    slide = content_slide(prs, "GAP NAMES")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Every gap has a name. Every defender owns a gap.",
                 font_size=18, color=MED_GRAY)

    data = [
        ["Gap", "Location"],
        ["A Gap", "Between Center and Guard"],
        ["B Gap", "Between Guard and Tackle"],
        ["C Gap", "Between Tackle and TE (or outside Tackle if no TE)"],
        ["D Gap", "Outside the TE"],
        ["E Gap", "Outside the Wing"],
        ["Alley", "Outside contain lane \u2014 where the force / contain player operates"],
    ]
    add_table(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.8),
              rows=7, cols=2, data=data,
              col_widths=[Inches(2), Inches(8.3)], font_size=18)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
                 "Gaps exist on both sides of the center. A/B/C/D field + A/B/C/D boundary.",
                 font_size=16, color=NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: CALL GRAMMAR
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "CALL GRAMMAR", "How We Communicate")

    # ── The Formula ──────────────────────────────────────────────────────
    slide = content_slide(prs, "THE FORMULA")

    formula_box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7),
                            Inches(1.4), FAINT_BLUE, COLUMBIA_BLUE)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7),
                                      Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    parts = [
        ("FRONT", NAVY),
        ("  +  ", DARK_GRAY),
        ("STUNT", RGBColor(0x8B, 0x00, 0x00)),
        ("  +  ", DARK_GRAY),
        ("BLITZ", RGBColor(0x00, 0x6B, 0x3C)),
        ("  +  ", DARK_GRAY),
        ("COVERAGE / TAGS", GOLD),
    ]
    for text, color in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = "Calibri"

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
                 "Examples:", font_size=22, color=NAVY, bold=True)

    examples = [
        [("SHADE", NAVY), (" + ", DARK_GRAY), ("NINJA", GOLD)],
        [("UNDER", NAVY), (" + ", DARK_GRAY),
         ("SLANT", RGBColor(0x8B, 0x00, 0x00)),
         (" + ", DARK_GRAY), ("ZEUS", GOLD)],
        [("GRIZZLY", NAVY), (" + ", DARK_GRAY),
         ("sWarM", RGBColor(0x00, 0x6B, 0x3C)),
         (" + ", DARK_GRAY), ("ZORRO", GOLD)],
        [("EYES", NAVY), (" + ", DARK_GRAY),
         ("JACKS", RGBColor(0x8B, 0x00, 0x00)),
         (" + ", DARK_GRAY), ("sWarM", RGBColor(0x00, 0x6B, 0x3C)),
         (" + ", DARK_GRAY), ("ZORRO", GOLD)],
    ]

    y_pos = Inches(3.9)
    for ex in examples:
        txBox = slide.shapes.add_textbox(Inches(1.5), y_pos, Inches(10),
                                          Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        for text, color in ex:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(22)
            run.font.bold = True if text.strip() not in ["+"] else False
            run.font.color.rgb = color
            run.font.name = "Calibri"
        y_pos += Inches(0.55)

    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
                 "Front/stunt/blitz calls signaled \u2014 coverage is verbal. Safeties echo and relay.",
                 font_size=16, color=MED_GRAY, bold=False)

    # ── Tempo Call ───────────────────────────────────────────────────────
    slide = content_slide(prs, "TEMPO CALL")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8),
                 "When the offense\u2019s tempo is faster than we can signal,\n"
                 "the defense defaults to the tempo call.",
                 font_size=20, color=DARK_GRAY)

    box = add_shape(slide, Inches(3), Inches(2.8), Inches(7.3), Inches(1.5),
                    NAVY, GOLD)
    add_text_box(slide, Inches(3), Inches(2.9), Inches(7.3), Inches(1.3),
                 "DEFAULT:  SHADE + NINJA", font_size=40, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    items = [
        "Set before each possession. Can change at any point during the game.",
        "Safeties must know the current tempo call at all times.",
        "Tempo call is our safety net \u2014 everyone knows exactly what to do.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(4.8), Inches(10.9), Inches(2),
                    items, font_size=18, line_spacing=1.5)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: SYSTEM OVERVIEW
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "THE SYSTEM", "What\u2019s in the Toolbox")

    # ── Fronts Menu ──────────────────────────────────────────────────────
    slide = content_slide(prs, "FRONT MENU  \u2014  4-DOWN")

    data = [
        ["Front", "Interior (A/T/N/E)", "Key Feature"],
        ["SHADE", "5 / 3 / 2i / 5", "Base front. 3-tech field, 2i boundary. Our home."],
        ["UNDER", "5 / 2i / 3 / 5", "Flipped interior. 3-tech to boundary."],
        ["EYES", "5 / 2i / 2i / 5", "Balanced interior. Both DTs inside shade."],
        ["WIDE", "5 / 3 / 3 / 5", "Both DTs outside shade. Closes B gaps, forces bounce."],
        ["DEUCES", "5 / 2 / 2 / 5", "Head-up guards. Read-and-react interior."],
        ["GRIZZLY", "4i / 2i / 2i / 4i", "Tight interior. Edges at 4i. Goal-line / short-yardage."],
        ["BOSS", "5 / 3(F) / 1(F) / 5", "Both DTs set field. Overload look."],
        ["BOSS UNDER", "5 / 1(B) / 3(B) / 5", "Both DTs set boundary. Overload look."],
    ]
    add_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5),
              rows=9, cols=3, data=data,
              col_widths=[Inches(2.2), Inches(2.8), Inches(7.5)], font_size=16)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
                 "TE SET rule overrides all fronts except GRIZZLY when TE is attached.\n"
                 "(F) = both DTs set to field. (B) = both DTs set to boundary.",
                 font_size=15, color=NAVY, bold=True)

    # ── 3-Down Packages ──────────────────────────────────────────────────
    slide = content_slide(prs, "3-DOWN PACKAGES")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Our odd-front package. Triggered by offensive personnel.",
                 font_size=18, color=MED_GRAY)

    data = [
        ["Package", "Interior (A/T/N)", "Key Feature"],
        ["MINT", "4i / 0 / 4i", "Base 3-down. Contain adjusts by coverage."],
        ["ACE", "4 / 0 / 4", "2-gap interior. Physical, read-and-react."],
        ["JET", "5 / 0 / 5", "Edges at 5. T is 2-gapping."],
        ["SLIP", "5 / 0 / 4i", "1-gap. Takes B out of run/pass conflict."],
    ]
    add_table(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.2),
              rows=5, cols=3, data=data,
              col_widths=[Inches(2.5), Inches(2.5), Inches(6.7)], font_size=18)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
                 "B = Field OLB.  E = Boundary OLB.  Strength = Field.",
                 font_size=16, color=NAVY, bold=True)

    # ── Stunts Overview ──────────────────────────────────────────────────
    slide = content_slide(prs, "STUNT MENU")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.6),
                 "Coordinated DL movement. In 4-down: stunts involve T+N. In 3-down: A joins.",
                 font_size=18, color=MED_GRAY)

    col1 = [
        ("SLANT ", "\u2014 Slant to field"),
        ("ANGLE ", "\u2014 Slant to boundary"),
        ("PINCH ", "\u2014 Pinch inside (A gaps / B gaps)"),
        ("JACKS ", "\u2014 Shoot out (B gaps / C gaps)"),
        ("VEER ", "\u2014 Slant to motion / to TE"),
        ("SPLIT ", "\u2014 BOSS / BOSS UNDER / GRIZZLY only"),
    ]
    col2 = [
        ("CRASH ", "\u2014 Everyone shoots inside"),
        ("ANCHOR ATTACK ", "\u2014 A washes OL inside"),
        ("EDGE ATTACK ", "\u2014 E washes OL inside"),
        ("ANCHOR RAVEN ", "\u2014 A shoots B gap"),
        ("EDGE RAVEN ", "\u2014 E shoots B gap"),
    ]

    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4),
                    col1, font_size=16, line_spacing=1.5)
    add_bullet_list(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4),
                    col2, font_size=16, line_spacing=1.5)

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
                 "Front-stunt legality matrix exists \u2014 not every stunt works with every front.",
                 font_size=15, color=NAVY, bold=True)

    # ── Pressures Overview ───────────────────────────────────────────────
    slide = content_slide(prs, "PRESSURE MENU")

    # Three columns
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.5),
                 "Singles", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.5),
                    ["Mike", "Will", "Bandit", "Dawg"],
                    font_size=18, line_spacing=1.4)

    add_text_box(slide, Inches(4.8), Inches(1.5), Inches(3.5), Inches(0.5),
                 "Combos", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(4.8), Inches(2.0), Inches(3.5), Inches(2.5),
                    ["sWarM (M+W)", "BooM (B+M)", "BoW (B+W)",
                     "MaD (M+D)", "Eat (M+W+B)"],
                    font_size=18, line_spacing=1.4)

    add_text_box(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(0.5),
                 "Packaged", font_size=20, color=NAVY, bold=True)
    add_bullet_list(slide, Inches(8.8), Inches(2.0), Inches(4), Inches(3.5),
                    [("Hammer ", "\u2014 B edge + Anchor Attack"),
                     ("Shave ", "\u2014 W edge + Edge Attack"),
                     ("Bandit Raven ", "\u2014 B blitzes C, A shoots B"),
                     ("Will Raven ", "\u2014 W blitzes C, E shoots B")],
                    font_size=18, line_spacing=1.4)

    # MUG note
    box = add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2),
                    FAINT_BLUE, COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.0),
                 "MUG: M and/or W walk up to the LOS at their gap. Same job, different picture.\n"
                 "Weekly menu locked by Friday. Game-plan specials installed as needed.",
                 font_size=16, color=DARK_GRAY)

    # ── Coverage Families ────────────────────────────────────────────────
    slide = content_slide(prs, "THE FOUR COVERAGE FAMILIES")

    y = Inches(1.5)
    box_h = Inches(1.15)
    gap = Inches(0.15)

    # NINJA
    add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE,
              COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(10), Inches(0.5),
                 "NINJA  \u2014  MOD / CLAMP / POACH", font_size=22,
                 color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Pattern-match quarters. Our base coverage. Most of our snaps. Two-high shell stays two-high.",
                 font_size=16, color=DARK_GRAY)

    # Cover 1
    y += box_h + gap
    add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE,
              COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 "COVER 1  (One Free Man)", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Oregon (FS = MOF) / Oklahoma (D = MOF) / Ohio (B = MOF). Man assignments + one free safety.",
                 font_size=16, color=DARK_GRAY)

    # Cover 0 / Z-Family
    y += box_h + gap
    add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE,
              COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(4), Inches(0.5),
                 "COVER 0 / Z-FAMILY", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Zero coverage. Zeus / Zorro / Zunnel / Zill / Zike. Run-first reads, then rush. No deep help.",
                 font_size=16, color=DARK_GRAY)

    # VIKING
    y += box_h + gap
    add_shape(slide, Inches(0.8), y, Inches(11.7), box_h, FAINT_BLUE,
              COLUMBIA_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.5), Inches(0.5),
                 "VIKING  (Cover 3)", font_size=22, color=NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.5),
                 "Spot-drop Cover 3. Safe, anti-explosive. Primary call-off after showing Zeus.",
                 font_size=16, color=DARK_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
                 "Zeus and Zorro are our identity. NINJA is the engine. Viking is the safety net.",
                 font_size=16, color=NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════════════
    # SECTION: HOW IT FITS TOGETHER — DIAGRAMS
    # ════════════════════════════════════════════════════════════════════════
    section_header_slide(prs, "HOW IT FITS TOGETHER",
                         "Same picture pre-snap. Different answers post-snap.")

    # ── SHADE + NINJA diagram ────────────────────────────────────────────
    slide = content_slide(prs, "SHADE + NINJA  \u2014  BASE LOOK")
    slide.shapes.add_picture(ninja_img, Inches(1.5), Inches(1.3),
                              Inches(10.3), Inches(5.8))

    # ── SHADE + ZEUS diagram ─────────────────────────────────────────────
    slide = content_slide(prs, "SHADE + ZEUS  \u2014  RUN-FIRST ZERO")
    slide.shapes.add_picture(zeus_img, Inches(1.5), Inches(1.3),
                              Inches(10.3), Inches(5.8))

    # ── SHADE + OREGON diagram ───────────────────────────────────────────
    slide = content_slide(prs, "SHADE + OREGON  \u2014  COVER 1 (ONE FREE)")
    slide.shapes.add_picture(oregon_img, Inches(1.5), Inches(1.3),
                              Inches(10.3), Inches(5.8))

    # ── Pressure look diagram ────────────────────────────────────────────
    slide = content_slide(prs, "SHADE + ANGLE + MIKE + ZEUS  \u2014  PRESSURE FROM DISGUISE")
    slide.shapes.add_picture(pressure_img, Inches(1.5), Inches(1.3),
                              Inches(10.3), Inches(5.8))

    # ── Three-look comparison diagram ────────────────────────────────────
    slide = content_slide(prs, "THREE SHELLS, THREE ANSWERS")
    slide.shapes.add_picture(comparison_img, Inches(0.5), Inches(1.3),
                              Inches(12.3), Inches(5.5))

    # ── System at a Glance ───────────────────────────────────────────────
    slide = content_slide(prs, "SYSTEM AT A GLANCE")

    data = [
        ["Category", "Count", "Detail"],
        ["4-Down Fronts", "8", "SHADE, UNDER, EYES, WIDE, DEUCES, GRIZZLY, BOSS, BOSS UNDER"],
        ["3-Down Packages", "4", "MINT, ACE, JET, SLIP"],
        ["Stunts", "11", "SLANT, ANGLE, PINCH, JACKS, SPLIT, CRASH, VEER, AA, EA, AR, ER"],
        ["Single Pressures", "4", "Mike, Will, Bandit, Dawg"],
        ["Combo Pressures", "5", "sWarM, BooM, BoW, MaD, Eat"],
        ["Packaged Pressures", "4", "Hammer, Shave, Bandit Raven, Will Raven"],
        ["Coverage Families", "4", "NINJA, Cover 1, Cover 0/Z, VIKING"],
        ["Z-Family Calls", "5", "Zeus, Zorro, Zunnel, Zill, Zike"],
        ["Alerts", "4", "BUMP, BANJO, EXCHANGE, UNDER"],
    ]
    add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.2),
              rows=10, cols=3, data=data,
              col_widths=[Inches(3.2), Inches(1.2), Inches(7.7)], font_size=16)

    add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
                 "Multiple looks. Simple, consistent rules. Layered install.",
                 font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDES 33-38 COMMENTED OUT — Camp install, in-season, 3-2-6, what's
    # coming section. Will be re-enabled once camp schedule is finalized.
    # ════════════════════════════════════════════════════════════════════════
    # section_header_slide(prs, "CAMP INSTALL PLAN", "The 10-Day Build")
    #
    # slide = content_slide(prs, "INSTALL OVERVIEW")
    # data = [
    #     ["Day", "Focus"],
    #     ["1", "SHADE front + ZEUS coverage + base run fits (IZ/OZ)"],
    #     ["2", "TE SET adjustment + ZEUS integration + Power/Counter"],
    #     ["3", "NINJA 2x2 (MOD/CLAMP) + BANJO/EXCHANGE alerts"],
    #     ["4", "NINJA 3x1 (POACH) + UNDER front + SLANT stunt"],
    #     ["5", "OREGON (Cover 1) + BUMP motion + checkpoint scrimmage"],
    #     ["6", "WIDE front + single pressures (Mike/Will/Bandit) + MUG"],
    #     ["7", "Combo pressures (sWarM/BooM/BoW) + ANGLE/PINCH stunts"],
    #     ["8", "EYES front + JACKS stunt + VIKING (Cover 3)"],
    #     ["9", "OHIO (Cover 1) + ZORRO/ZUNNEL + UNDER alert"],
    #     ["10", "MINT package + full integration scrimmage + situational review"],
    # ]
    # add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4),
    #           rows=11, cols=2, data=data,
    #           col_widths=[Inches(1.2), Inches(10.9)], font_size=16)
    # add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.4),
    #              "Zeus installs Day 1. NINJA Day 3. Everything layers from there.",
    #              font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    #
    # slide = content_slide(prs, "IN-SEASON ADDITIONS")
    # data = [
    #     ["Timing", "Add"],
    #     ["Weeks 1-2", "Refine camp installs. Add GRIZZLY if facing power/counter."],
    #     ["Weeks 3-4", "OKLAHOMA (Cover 1). ACE/JET packages. CRASH stunt."],
    #     ["Weeks 5-6", "DEUCES. Ravens. ZILL/ZIKE. Hammer."],
    #     ["Weeks 7-8", "BOSS/BOSS UNDER. SPLIT. Cobra + Zunnel. Eat."],
    #     ["Weeks 9+", "SLIP package. VIKING tags (SEAM, PUSH, CROSS, SCREEN)."],
    #     ["Playoffs", "Game-plan specials (Dawg, MaD, Shave, Freebird) as needed."],
    # ]
    # add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.2),
    #           rows=7, cols=2, data=data,
    #           col_widths=[Inches(2.5), Inches(9.2)], font_size=18)
    # add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
    #              "Nothing gets added until what's installed is clean. Quality over quantity.",
    #              font_size=16, color=NAVY, bold=True)
    #
    # section_header_slide(prs, "WHAT'S COMING", "")
    #
    # slide = content_slide(prs, "ON THE HORIZON: 3-2-6 SHELL")
    # add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
    #              "We are developing a 3-2-6 package in collaboration with Coach Arndt.",
    #              font_size=22, color=DARK_GRAY)
    # items = [
    #     ("What it is: ", "3 DL, 2 LBs, 6 DBs. An additional DB replaces a DL or LB."),
    #     ("Why: ", "Gives us another tool against spread-heavy and empty formations."),
    #     ("Status: ", "Still in development. Will be integrated once finalized."),
    #     ("How it fits: ", "Same coverage families, same call grammar. Different personnel."),
    # ]
    # add_bullet_list(slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(3),
    #                 items, font_size=20, line_spacing=1.5)
    # add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
    #              "The system is built to absorb new packages without changing the rules.",
    #              font_size=16, color=NAVY, bold=True)

    # ── Next Steps ───────────────────────────────────────────────────────
    slide = content_slide(prs, "NEXT STEPS")

    items = [
        ("Playbook ", "\u2014 Full document (27+ sections) will be shared via Google Drive after this meeting."),
        ("Unit meetings ", "\u2014 DL, LBs, DBs will each get deep-dives on their responsibilities."),
        ("Diagram documents ", "\u2014 Visual references for every front, coverage, and pressure."),
        ("Quick-reference cards ", "\u2014 M/W gap fits, NINJA teaching summary, coverage families, Bandit assignment."),
        ("Camp install plan ", "\u2014 Coming once camp schedule is finalized."),
        ("3-2-6 shell ", "\u2014 In development with Coach J. Arndt. Same rules, different personnel package."),
        ("This presentation ", "\u2014 Will be on the Drive after the meeting."),
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5),
                    items, font_size=20, line_spacing=1.5)

    # ── Closing ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_gold_bar(slide, Inches(2.0))
    add_gold_bar(slide, Inches(5.5))

    # Philosophy callback
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3),
                                      Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(42)
    lines = [
        "NO EXPLOSIVES.  NO FREE YARDS.",
        "MAKE THEM EARN IT."
    ]
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = Pt(42)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(40)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"

    add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
                 "River Valley Vikings | 2026", font_size=24, color=GOLD,
                 bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
                 "Questions?", font_size=28, color=COLUMBIA_BLUE,
                 bold=False, alignment=PP_ALIGN.CENTER)

    # ── Save ────────────────────────────────────────────────────────────
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir,
                            "Meeting_01_V2_Introduction_to_the_Defense.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
